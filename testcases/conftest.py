"""Shared pytest configuration and fixture-generation helpers for docintel_kit tests.

This module builds a `fixtures/` directory (generated once per test session,
cached on disk) containing deliberately tricky files across every format the
library claims to support:

- PDFs: multi-page text, unicode/CJK/RTL text, an image-only (scanned) page,
  an empty page, a corrupted/truncated file.
- DOCX: paragraphs + a table + unicode text.
- PPTX: multiple slides with text and a table.
- HTML: malformed markup, unicode, script/style stripping.
- Images: a clean synthetic "form" image, a rotated version, a blank
  (textless) image, and a corrupted image file.
- XLSX: multiple sheets, NaN/blank cells, mixed types.
- CSV/TSV: quoted commas, unicode, ragged rows.
- An unsupported file type (`.xyz`) to exercise error handling.

Fixtures are session-scoped and written under `testcases/fixtures/`. They are
regenerated on every test run to avoid staleness across docintel_kit changes;
`fixtures/` is safe to delete at any time.
"""

from __future__ import annotations

import os
import struct
from pathlib import Path

import pytest

FIXTURES_DIR = Path(__file__).parent / "fixtures"


def _ensure_dir() -> Path:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    return FIXTURES_DIR


# --------------------------------------------------------------------------
# PDF fixtures
# --------------------------------------------------------------------------


_UNICODE_FONT_CANDIDATES = [
    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
    "/Library/Fonts/Arial Unicode.ttf",
]


def _register_unicode_font() -> str:
    """Register a TTF with broad Unicode coverage (CJK/Greek/Cyrillic) for
    reportlab and return its font name. Falls back to Helvetica (which lacks
    these glyphs) if no suitable font is found on the system — callers should
    treat the fallback as "ASCII/Latin-1 only" for test purposes.
    """
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    for font_path in _UNICODE_FONT_CANDIDATES:
        if Path(font_path).exists():
            pdfmetrics.registerFont(TTFont("UnicodeFont", font_path))
            return "UnicodeFont"
    return "Helvetica"


def _make_text_pdf(path: Path) -> None:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    font_name = _register_unicode_font()
    c = canvas.Canvas(str(path), pagesize=letter)
    for page_num in range(3):
        c.setFont(font_name, 14)
        c.drawString(72, 720, f"Page {page_num + 1} of a synthetic report")
        c.setFont(font_name, 11)
        c.drawString(72, 690, "Quarterly revenue increased by 12% year over year.")
        c.drawString(72, 670, "Unicode check: café, naïve, 日本語, Ελληνικά, Русский текст.")
        # Emoji glyphs are excluded here: even broad-coverage TTFs like Arial
        # Unicode.ttf don't include color emoji, so this is intentionally not
        # asserted on in tests. Emoji handling for a *text-layer* PDF is a
        # font/rendering concern, not a docintel_kit parsing concern.
        c.showPage()
    c.save()


def _make_empty_pdf(path: Path) -> None:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(path), pagesize=letter)
    c.showPage()  # a page with literally nothing drawn on it
    c.save()


def _make_image_only_pdf(path: Path) -> None:
    """A PDF whose only page is a rasterized image (no text layer at all) —
    simulates a scanned document, forcing OCR fallback."""
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    from PIL import Image, ImageDraw, ImageFont

    img = Image.new("RGB", (1200, 1600), "white")
    draw = ImageDraw.Draw(img)
    title_font = _load_ocr_font(42)
    body_font = _load_ocr_font(34)
    draw.text((80, 80), "SCANNED INVOICE", fill="black", font=title_font)
    draw.text((80, 170), "Invoice Number: INV-90210", fill="black", font=body_font)
    draw.text((80, 230), "Total Due: $4,521.00", fill="black", font=body_font)
    draw.rectangle([80, 300, 1100, 700], outline="black", width=3)
    img_path = path.with_suffix(".scan.png")
    img.save(img_path)

    c = canvas.Canvas(str(path), pagesize=letter)
    c.drawImage(str(img_path), 0, 0, width=letter[0], height=letter[1])
    c.showPage()
    c.save()
    img_path.unlink(missing_ok=True)


def _make_table_pdf(path: Path) -> None:
    """A single-page PDF with a real ruled (bordered) table, drawn with
    vector lines so camelot's default 'lattice' flavor can detect it."""
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(path), pagesize=letter)
    c.setFont("Helvetica", 16)
    c.drawString(72, 740, "Invoice INV-55219")

    rows = [
        ["Item", "Qty", "Unit Price", "Total"],
        ["Widget A", "12", "$5.00", "$60.00"],
        ["Widget B", "3", "$19.99", "$59.97"],
        ["Gizmo", "1", "$250.00", "$250.00"],
    ]
    col_widths = [180, 60, 100, 100]
    row_height = 24
    x0, y0 = 72, 650
    table_width = sum(col_widths)
    table_height = row_height * len(rows)

    # Horizontal lines
    for r in range(len(rows) + 1):
        y = y0 - r * row_height
        c.line(x0, y, x0 + table_width, y)
    # Vertical lines
    x = x0
    for w in [0] + col_widths:
        x += w
        c.line(x, y0, x, y0 - table_height)

    c.setFont("Helvetica", 10)
    for r, row in enumerate(rows):
        x = x0 + 4
        y = y0 - r * row_height - row_height + 8
        for ci, cell in enumerate(row):
            c.drawString(x, y, str(cell))
            x += col_widths[ci]

    c.showPage()
    c.save()


def _make_corrupted_pdf(path: Path) -> None:
    path.write_bytes(b"%PDF-1.7\nThis is not a valid PDF body at all.\n%%EOF-broken")


def _make_truncated_pdf(source_path: Path, dest_path: Path) -> None:
    data = source_path.read_bytes()
    dest_path.write_bytes(data[: len(data) // 3])


# --------------------------------------------------------------------------
# DOCX / PPTX fixtures
# --------------------------------------------------------------------------


def _make_docx(path: Path) -> None:
    import docx

    doc = docx.Document()
    doc.add_heading("Contract Agreement", level=1)
    doc.add_paragraph("This agreement is made between Party A and Party B.")
    doc.add_paragraph("Unicode check: café, naïve, 日本語, Ελληνικά.")
    table = doc.add_table(rows=3, cols=3)
    headers = ["Item", "Qty", "Price"]
    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = h
    table.rows[1].cells[0].text = "Widget"
    table.rows[1].cells[1].text = "10"
    table.rows[1].cells[2].text = "$5.00"
    table.rows[2].cells[0].text = "Gadget"
    table.rows[2].cells[1].text = "3"
    table.rows[2].cells[2].text = "$19.99"
    doc.save(str(path))


def _make_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Q3 Results"
    slide1.placeholders[1].text = "Revenue grew 12% QoQ.\nUnicode: café, 日本語."

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    rows, cols = 2, 2
    table_shape = slide2.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(4), Inches(2))
    table = table_shape.table
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Growth"
    table.cell(1, 0).text = "APAC"
    table.cell(1, 1).text = "18%"

    prs.save(str(path))


# --------------------------------------------------------------------------
# HTML fixtures
# --------------------------------------------------------------------------


def _make_malformed_html(path: Path) -> None:
    html = """<html><head><title>Broken</title>
<style>body { color: red; }</style>
<script>alert('should be stripped');</script>
</head><body>
<h1>Unclosed heading
<p>Paragraph with <b>bold and <i>nested unclosed tags
<div>Unicode: café, naïve, 日本語, Ελληνικά, emoji 🚀</div>
<table><tr><td>Cell1<td>Cell2</tr></table>
</body>
"""
    path.write_text(html, encoding="utf-8")


# --------------------------------------------------------------------------
# Image fixtures
# --------------------------------------------------------------------------


_LATIN_FONT_CANDIDATES = [
    "/System/Library/Fonts/Supplemental/Arial.ttf",
    "/System/Library/Fonts/Helvetica.ttc",
]


def _load_ocr_font(size: int):
    """Load a scalable TTF for OCR test images.

    PIL's ImageDraw default font is a tiny fixed-size bitmap font that
    Tesseract struggles to recognize accurately, especially after rotation.
    Using a real, appropriately-sized TTF makes these fixtures representative
    of realistic scanned documents rather than an artifact of the test setup.
    """
    from PIL import ImageFont

    for font_path in _LATIN_FONT_CANDIDATES:
        if Path(font_path).exists():
            try:
                return ImageFont.truetype(font_path, size)
            except OSError:
                continue
    return ImageFont.load_default()


def _make_form_image(path: Path) -> None:
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (900, 500), "white")
    draw = ImageDraw.Draw(img)
    title_font = _load_ocr_font(32)
    body_font = _load_ocr_font(26)
    draw.text((40, 30), "PURCHASE ORDER", fill="black", font=title_font)
    draw.text((40, 100), "PO Number: PO-778812", fill="black", font=body_font)
    draw.text((40, 150), "Date: 2026-01-15", fill="black", font=body_font)
    draw.text((40, 200), "Total: $1,250.75", fill="black", font=body_font)
    draw.rectangle([30, 20, 870, 480], outline="black", width=2)
    img.save(path)


def _make_rotated_image(source_path: Path, dest_path: Path) -> None:
    from PIL import Image

    img = Image.open(source_path)
    rotated = img.rotate(7, expand=True, fillcolor="white")
    rotated.save(dest_path)


def _make_blank_image(path: Path) -> None:
    from PIL import Image

    img = Image.new("RGB", (400, 400), "white")
    img.save(path)


def _make_corrupted_image(path: Path) -> None:
    path.write_bytes(b"\x89PNG\r\n\x1a\n" + os.urandom(64))


# --------------------------------------------------------------------------
# Spreadsheet fixtures
# --------------------------------------------------------------------------


def _make_xlsx(path: Path) -> None:
    import pandas as pd

    with pd.ExcelWriter(path, engine="openpyxl") as writer:
        df1 = pd.DataFrame(
            {
                "Product": ["Widget", "Gadget", "Gizmo", None],
                "Units Sold": [120, 85, None, 40],
                "Revenue": [1200.50, 850.0, 300.25, float("nan")],
            }
        )
        df1.to_excel(writer, sheet_name="Sales", index=False)

        df2 = pd.DataFrame(
            {
                "Region": ["APAC", "EMEA", "AMER"],
                "Growth %": [18.2, -3.5, 7.0],
                "Notes": ["café", "naïve résumé", "日本語ノート"],
            }
        )
        df2.to_excel(writer, sheet_name="Regions", index=False)


def _make_csv(path: Path) -> None:
    content = (
        'name,description,amount\n'
        '"Smith, John","A ""quoted"" description with, commas",1234.56\n'
        "Café Owner,Unicode naïve test 日本語,99.99\n"
        "Ragged Row,only two fields\n"
    )
    path.write_text(content, encoding="utf-8")


def _make_tsv(path: Path) -> None:
    content = "name\tamount\nAlice\t100\nBöb\t250\n"
    path.write_text(content, encoding="utf-8")


# --------------------------------------------------------------------------
# Misc
# --------------------------------------------------------------------------


def _make_unsupported_file(path: Path) -> None:
    path.write_bytes(b"\x00\x01\x02just some random binary junk\x03\x04")


@pytest.fixture(scope="session")
def fixtures_dir() -> Path:
    """Build (or reuse) the full set of high-difficulty test fixtures, return their dir."""
    d = _ensure_dir()

    text_pdf = d / "text_report.pdf"
    _make_text_pdf(text_pdf)
    _make_empty_pdf(d / "empty_page.pdf")
    _make_image_only_pdf(d / "scanned_invoice.pdf")
    _make_table_pdf(d / "invoice_table.pdf")
    _make_corrupted_pdf(d / "corrupted.pdf")
    _make_truncated_pdf(text_pdf, d / "truncated.pdf")

    _make_docx(d / "contract.docx")
    _make_pptx(d / "slides.pptx")

    _make_malformed_html(d / "malformed.html")

    form_image = d / "form.png"
    _make_form_image(form_image)
    _make_rotated_image(form_image, d / "form_rotated.png")
    _make_blank_image(d / "blank.png")
    _make_corrupted_image(d / "corrupted.png")

    _make_xlsx(d / "workbook.xlsx")
    _make_csv(d / "tricky.csv")
    _make_tsv(d / "tricky.tsv")

    _make_unsupported_file(d / "unknown.xyz")

    return d
