"""Document parsing & normalization.

Converts PDF, DOCX, PPTX, HTML, and image files into a normalized
:class:`~docintel_kit.types.ParseResult` (per-page plain text plus metadata),
using local, dependency-light parsers:

- PDF: `pdfplumber <https://github.com/jsvine/pdfplumber>`_
- DOCX: `python-docx <https://python-docx.readthedocs.io/>`_
- PPTX: `python-pptx <https://python-pptx.readthedocs.io/>`_
- HTML: `BeautifulSoup <https://www.crummy.com/software/BeautifulSoup/>`_
- Images: no native text layer; returns an empty page and a warning pointing
  callers to :func:`docintel_kit.ocr.run_ocr`.

Backend selection is pluggable via :class:`BaseParserBackend` so alternative
implementations (e.g. a commercial parser, or a faster native binding) can be
swapped in without touching call sites.
"""

from __future__ import annotations

import hashlib
import mimetypes
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Optional, Union

from .types import Document, Page, ParseResult

__all__ = ["BaseParserBackend", "DefaultParserBackend", "parse_document"]

# MIME types that the default backend knows how to route.
_PDF_MIME = "application/pdf"
_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_HTML_MIMES = {"text/html", "application/xhtml+xml"}
_IMAGE_PREFIX = "image/"

_EXTENSION_MIME_OVERRIDES = {
    ".pdf": _PDF_MIME,
    ".docx": _DOCX_MIME,
    ".pptx": _PPTX_MIME,
    ".htm": "text/html",
    ".html": "text/html",
}


def _make_document_id(input_: Union[str, bytes], source_path: Optional[str]) -> str:
    """Derive a stable identifier: the file path if we have one, else a content hash."""
    if source_path:
        return source_path
    data = input_ if isinstance(input_, bytes) else input_.encode("utf-8")
    return hashlib.sha256(data).hexdigest()[:32]


def _guess_mime_type(input_: Union[str, bytes], mime_type: Optional[str]) -> str:
    """Resolve a MIME type from an explicit hint, file extension, or content sniffing."""
    if mime_type:
        return mime_type

    if isinstance(input_, str):
        suffix = Path(input_).suffix.lower()
        if suffix in _EXTENSION_MIME_OVERRIDES:
            return _EXTENSION_MIME_OVERRIDES[suffix]
        guessed, _ = mimetypes.guess_type(input_)
        if guessed:
            return guessed
        raise ValueError(
            f"Could not infer mime_type for '{input_}'. Pass mime_type explicitly."
        )

    # Raw bytes: sniff common magic numbers.
    if input_.startswith(b"%PDF"):
        return _PDF_MIME
    if input_[:4] == b"PK\x03\x04":
        # Zip-based Office format; DOCX vs PPTX can't be told apart from the
        # header alone, default to DOCX and let callers override via mime_type.
        return _DOCX_MIME
    if input_.lstrip()[:1] in (b"<",):
        return "text/html"
    raise ValueError(
        "Could not infer mime_type from raw bytes. Pass mime_type explicitly."
    )


class BaseParserBackend(ABC):
    """Interface for a document parsing backend.

    Implement this to plug in an alternative parsing engine while keeping
    :func:`parse_document` and its return type (:class:`ParseResult`) stable
    for callers.
    """

    name: str = "base"

    @abstractmethod
    def supports(self, mime_type: str) -> bool:
        """Return True if this backend can handle the given MIME type."""
        raise NotImplementedError

    @abstractmethod
    def parse(
        self,
        input_: Union[str, bytes],
        mime_type: str,
        document: Document,
    ) -> ParseResult:
        """Parse ``input_`` (a path or raw bytes) into a :class:`ParseResult`."""
        raise NotImplementedError


class DefaultParserBackend(BaseParserBackend):
    """Local parser backend built on pdfplumber, python-docx, python-pptx, and BeautifulSoup."""

    name = "default"

    def supports(self, mime_type: str) -> bool:
        return (
            mime_type == _PDF_MIME
            or mime_type == _DOCX_MIME
            or mime_type == _PPTX_MIME
            or mime_type in _HTML_MIMES
            or mime_type.startswith(_IMAGE_PREFIX)
        )

    def parse(
        self,
        input_: Union[str, bytes],
        mime_type: str,
        document: Document,
    ) -> ParseResult:
        if mime_type == _PDF_MIME:
            return self._parse_pdf(input_, document)
        if mime_type == _DOCX_MIME:
            return self._parse_docx(input_, document)
        if mime_type == _PPTX_MIME:
            return self._parse_pptx(input_, document)
        if mime_type in _HTML_MIMES:
            return self._parse_html(input_, document)
        if mime_type.startswith(_IMAGE_PREFIX):
            return self._parse_image(document)
        raise ValueError(f"DefaultParserBackend cannot handle mime_type={mime_type!r}")

    # -- format-specific implementations ---------------------------------

    def _parse_pdf(self, input_: Union[str, bytes], document: Document) -> ParseResult:
        import io

        import pdfplumber

        source = io.BytesIO(input_) if isinstance(input_, bytes) else input_
        pages: list[Page] = []
        warnings: list[str] = []
        with pdfplumber.open(source) as pdf:
            for index, pdf_page in enumerate(pdf.pages):
                text = pdf_page.extract_text() or ""
                if not text.strip():
                    warnings.append(
                        f"Page {index} has no extractable text layer; it may be a "
                        "scanned image. Consider docintel_kit.ocr.run_ocr()."
                    )
                pages.append(
                    Page(
                        index=index,
                        width=float(pdf_page.width),
                        height=float(pdf_page.height),
                        text=text,
                    )
                )
        document.page_count = len(pages)
        return ParseResult(document=document, pages=pages, backend=self.name, warnings=warnings)

    def _parse_docx(self, input_: Union[str, bytes], document: Document) -> ParseResult:
        import io

        import docx

        source = io.BytesIO(input_) if isinstance(input_, bytes) else input_
        doc = docx.Document(source)
        paragraphs = [p.text for p in doc.paragraphs]
        table_lines: list[str] = []
        for table in doc.tables:
            for row in table.rows:
                table_lines.append(" | ".join(cell.text for cell in row.cells))
        text = "\n".join(paragraphs + table_lines)
        # DOCX has no fixed page geometry until rendered/paginated, so we treat
        # the whole document as a single logical page.
        page = Page(index=0, text=text)
        document.page_count = 1
        return ParseResult(document=document, pages=[page], backend=self.name)

    def _parse_pptx(self, input_: Union[str, bytes], document: Document) -> ParseResult:
        import io

        from pptx import Presentation

        source = io.BytesIO(input_) if isinstance(input_, bytes) else input_
        presentation = Presentation(source)
        pages: list[Page] = []
        for index, slide in enumerate(presentation.slides):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        run_text = "".join(run.text for run in paragraph.runs)
                        if run_text:
                            texts.append(run_text)
                if shape.has_table:
                    for row in shape.table.rows:
                        texts.append(" | ".join(cell.text for cell in row.cells))
            pages.append(
                Page(
                    index=index,
                    width=float(presentation.slide_width),
                    height=float(presentation.slide_height),
                    text="\n".join(texts),
                )
            )
        document.page_count = len(pages)
        return ParseResult(document=document, pages=pages, backend=self.name)

    def _parse_html(self, input_: Union[str, bytes], document: Document) -> ParseResult:
        from bs4 import BeautifulSoup

        if isinstance(input_, bytes):
            html = self._decode_html_bytes(input_)
        else:
            path = Path(input_)
            if path.exists():
                html = self._decode_html_bytes(path.read_bytes())
            else:
                html = input_

        soup = BeautifulSoup(html, "lxml")
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        page = Page(index=0, text=text)
        document.page_count = 1
        return ParseResult(document=document, pages=[page], backend=self.name)

    @staticmethod
    def _decode_html_bytes(data: bytes) -> str:
        """Decode raw bytes into an HTML string, rejecting obviously-binary content.

        Guards against the common misuse-case of an explicit (wrong)
        ``mime_type="text/html"`` being passed for a binary file (PDF, ZIP-based
        Office formats, ...): rather than silently decoding garbage bytes via
        ``errors="replace"`` and handing nonsense to BeautifulSoup, this
        recognizes known binary magic numbers and raises immediately. Genuine
        non-UTF-8 HTML (e.g. legacy Latin-1 pages) still decodes successfully.
        """
        _BINARY_MAGIC_PREFIXES = (
            b"%PDF",  # PDF
            b"PK\x03\x04",  # ZIP / DOCX / PPTX / XLSX
            b"\xd0\xcf\x11\xe0",  # legacy OLE (.doc/.xls/.ppt)
            b"\x89PNG",
            b"\xff\xd8\xff",  # JPEG
            b"GIF8",
        )
        if data.startswith(_BINARY_MAGIC_PREFIXES):
            raise ValueError(
                "Input looks like a binary file (PDF/Office/image), not HTML. "
                "Check that mime_type is correct for this input."
            )
        try:
            return data.decode("utf-8")
        except UnicodeDecodeError:
            return data.decode("latin-1")

    def _parse_image(self, document: Document) -> ParseResult:
        page = Page(index=0, text="")
        document.page_count = 1
        return ParseResult(
            document=document,
            pages=[page],
            backend=self.name,
            warnings=[
                "Images have no native text layer; call docintel_kit.ocr.run_ocr() "
                "to extract text via OCR."
            ],
        )


# Registry of available backends, keyed by name. Populated with the default
# backend; additional backends can be registered by callers via
# `register_parser_backend`.
_BACKENDS: dict[str, BaseParserBackend] = {"default": DefaultParserBackend()}


def register_parser_backend(backend: BaseParserBackend) -> None:
    """Register a custom :class:`BaseParserBackend` under its ``name``."""
    _BACKENDS[backend.name] = backend


def parse_document(
    input: Union[str, bytes],
    mime_type: Optional[str] = None,
    backend: str = "default",
) -> ParseResult:
    """Parse a document into normalized per-page text.

    Args:
        input: Either a filesystem path (str) or the raw file contents (bytes).
        mime_type: Explicit MIME type. If omitted, it is inferred from the
            file extension (path input) or from content sniffing (bytes input).
        backend: Name of a registered :class:`BaseParserBackend` to use.
            Defaults to the built-in local backend covering PDF, DOCX, PPTX,
            HTML, and images.

    Returns:
        A :class:`ParseResult` with per-page text and any parsing warnings
        (e.g. a PDF page with no extractable text, suggesting OCR).

    Raises:
        ValueError: if ``mime_type`` cannot be inferred, or no registered
            backend supports the resolved MIME type.
        KeyError: if ``backend`` is not a registered backend name.
    """
    resolved_mime = _guess_mime_type(input, mime_type)
    source_path = input if isinstance(input, str) else None
    document = Document(
        id=_make_document_id(input, source_path),
        source_path=source_path,
        mime_type=resolved_mime,
    )

    if backend not in _BACKENDS:
        raise KeyError(
            f"Unknown parser backend '{backend}'. Registered backends: "
            f"{sorted(_BACKENDS)}"
        )
    selected = _BACKENDS[backend]
    if not selected.supports(resolved_mime):
        raise ValueError(
            f"Backend '{backend}' does not support mime_type={resolved_mime!r}"
        )
    return selected.parse(input, resolved_mime, document)
